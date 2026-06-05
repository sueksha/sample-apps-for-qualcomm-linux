#!/usr/bin/env python3
"""
Build a cleaner, data-driven executive deck for Qualcomm GenAI Studio on IQ-9075 EVK
versus NVIDIA Jetson edge platforms.

Design goals:
- Fewer words per slide
- Clear public vs internal vs roadmap separation
- More chart-driven / matrix-driven communication
- Safer spacing to avoid text collisions
"""

from __future__ import annotations

from pathlib import Path
from typing import Iterable, Sequence

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "artifacts"
OUT_DIR.mkdir(exist_ok=True)

STAMP = "2026-04-28"
OUT_PPTX = OUT_DIR / f"Qualcomm_GenAI_Studio_IQ9075_vs_Jetson_Strategy_EXEC_v2_{STAMP}.pptx"
OUT_NOTES = OUT_DIR / f"Qualcomm_GenAI_Studio_IQ9075_vs_Jetson_Strategy_EXEC_v2_{STAMP}_notes.md"
OUT_SOURCES = OUT_DIR / f"Qualcomm_GenAI_Studio_IQ9075_vs_Jetson_Strategy_EXEC_v2_{STAMP}_sources.md"


class C:
    BG = RGBColor(245, 243, 237)
    INK = RGBColor(22, 29, 45)
    MUTED = RGBColor(92, 101, 122)
    DEEP = RGBColor(33, 53, 85)
    CYAN = RGBColor(0, 161, 198)
    GOLD = RGBColor(214, 164, 32)
    GREEN = RGBColor(86, 144, 78)
    RED = RGBColor(191, 84, 71)
    PLUM = RGBColor(123, 88, 156)
    PANEL = RGBColor(255, 255, 255)
    PANEL_ALT = RGBColor(236, 239, 245)
    BORDER = RGBColor(205, 211, 224)
    JETSON = RGBColor(118, 185, 0)
    QCOM = RGBColor(0, 103, 177)
    AMBER_BG = RGBColor(255, 246, 214)
    BLUE_BG = RGBColor(226, 241, 248)
    PURPLE_BG = RGBColor(240, 233, 248)


PUBLIC = ("PUBLIC", C.CYAN, C.BLUE_BG)
INTERNAL = ("ENGINEERING EVIDENCE", C.GOLD, C.AMBER_BG)
ROADMAP = ("ROADMAP", C.PLUM, C.PURPLE_BG)

SW = 13.333
SH = 7.5


SOURCE_MAP = {
    "Q1": "Qualcomm GenAI Studio docs (updated Feb 23, 2026) - https://docs.qualcomm.com/doc/80-90441-15/topic/develop-genai-app-with-genai-studio.html",
    "Q2": "Public GenAI Studio GitHub README - https://github.com/quic/sample-apps-for-qualcomm-linux/blob/main/GenAI-Solutions/GenAI-Studio/README.md",
    "Q3": "Qualcomm IQ-9075 product page - https://www.qualcomm.com/internet-of-things/products/iq9-series/iq-9075",
    "Q4": "Qualcomm IQ-9075 EVK developer page - https://www.qualcomm.com/developer/hardware/qualcomm-iq-9075-evaluation-kit-evk",
    "Q5": "Qualcomm IQ-9075 EVK brief PDF - https://www.qualcomm.com/content/dam/qcomm-martech/dm-assets/documents/Qualcomm-Dragonwing-IQ-9075-EVK-Brief.pdf",
    "H1": "Qualcomm AI Hub home - https://aihub.qualcomm.com/",
    "H2": "Qualcomm AI Hub get started - https://aihub.qualcomm.com/get-started",
    "H3": "Qualcomm AI Hub Apps - https://aihub.qualcomm.com/apps",
    "H4": "Qualcomm AI Hub Workbench docs - https://workbench.aihub.qualcomm.com/docs/",
    "N1": "NVIDIA Jetson Orin Nano Super page - https://www.nvidia.com/en-us/autonomous-machines/embedded-systems/jetson-orin/nano-super-developer-kit/",
    "N2": "NVIDIA JetPack SDK downloads and notes - https://developer.nvidia.com/embedded/jetpack/downloads",
    "N3": "NVIDIA Jetson tutorials / Generative AI Lab - https://developer.nvidia.com/embedded/learn/tutorials",
    "N4": "NVIDIA Jetson software getting started - https://developer.nvidia.com/embedded/learn/getting-started-jetson",
    "N5": "NVIDIA Jetson tutorials / GenAI Lab - https://developer.nvidia.com/embedded/learn/tutorials",
    "N6": "NVIDIA Jetson edge AI blog (Dec 11, 2025) - https://developer.nvidia.com/blog/getting-started-with-edge-ai-on-nvidia-jetson-llms-vlms-and-foundation-models-for-robotics/",
    "N7": "Isaac ROS landing / getting started - https://nvidia-isaac-ros.github.io/",
    "M1": "MarketsandMarkets edge AI hardware market (Jun 2025) - https://www.marketsandmarkets.com/Market-Reports/edge-ai-hardware-market-158498281.html",
    "M2": "Global Market Insights edge AI hardware market (Jul 2025) - https://www.gminsights.com/industry-analysis/edge-ai-hardware-market",
    "M3": "Edge AI & Vision / SHD report release (Jun 2025) - https://www.edge-ai-vision.com/2025/06/the-shd-group-releases-new-edge-ai-processor-and-ecosystem-report/",
    "I1": "Local repo: README.md",
    "I2": "Local repo: core-services/orchestrator/README.md",
    "I3": "Local repo: docs/API_CONTRACTS.md",
    "I4": "Local repo: core-services/image-to-text/README.md",
    "I5": "Local repo: core-services/speech-to-text/README.md",
    "I6": "Local repo: core-services/orchestrator/app/openai.py",
}


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)

    notes: list[str] = []

    slide_title(prs, notes)
    slide_why_now(prs, notes)
    slide_public_proof(prs, notes)
    slide_hardware(prs, notes)
    slide_internal(prs, notes)
    slide_aihub_growth(prs, notes)
    slide_unified_endpoints(prs, notes)
    slide_dev_value(prs, notes)
    slide_customer_value(prs, notes)
    slide_jetson_advantage(prs, notes)
    slide_qcom_gaps(prs, notes)
    slide_gap_closure(prs, notes)
    slide_competitive_snapshot(prs, notes)
    slide_priorities(prs, notes)
    slide_appendix(prs, notes)

    prs.save(str(OUT_PPTX))
    OUT_NOTES.write_text("\n\n".join(notes), encoding="utf-8")
    OUT_SOURCES.write_text(render_sources(), encoding="utf-8")


def render_sources() -> str:
    lines = [
        "# Qualcomm GenAI Studio vs Jetson executive deck sources",
        "",
        f"Generated: {STAMP}",
        "",
        "## Public sources",
    ]
    for key in [k for k in SOURCE_MAP if not k.startswith("I")]:
        lines.append(f"- `{key}` {SOURCE_MAP[key]}")
    lines.extend(["", "## Internal engineering evidence", ""])
    for key in [k for k in SOURCE_MAP if k.startswith("I")]:
        lines.append(f"- `{key}` {SOURCE_MAP[key]}")
    lines.extend(
        [
            "",
            "## Truth model used in the deck",
            "",
            "- `PUBLIC`: backed by public Qualcomm / NVIDIA / market sources.",
            "- `ENGINEERING EVIDENCE`: backed by the current local repo and should be validated with owners before external use.",
            "- `ROADMAP`: recommendation / next-step narrative; not a public commitment.",
        ]
    )
    return "\n".join(lines)


def add_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = C.BG
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(SW), Inches(0.23)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = C.DEEP
    band.line.fill.background()


def add_tag(slide, text: str, color: RGBColor, bg: RGBColor, x: float, y: float, w: float = 1.75) -> None:
    shp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.34)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = bg
    shp.line.color.rgb = color
    shp.line.width = Pt(1.2)
    tf = shp.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = color
    r.font.name = "Aptos"


def add_title(slide, title: str, subtitle: str, tag: tuple[str, RGBColor, RGBColor], slide_no: int, refs: str) -> None:
    add_bg(slide)
    add_tag(slide, tag[0], tag[1], tag[2], 0.68, 0.42)
    add_text(slide, 0.68, 0.82, 11.0, 0.78, title, 24, C.INK, True, "Aptos Display")
    add_text(slide, 0.68, 1.42, 11.0, 0.24, subtitle, 11, C.MUTED, False, "Aptos")
    add_footer(slide, refs, slide_no)


def add_footer(slide, refs: str, slide_no: int) -> None:
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.6), Inches(6.86), Inches(12.1), Inches(0.01)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = C.BORDER
    line.line.fill.background()
    add_text(slide, 0.68, 6.9, 11.1, 0.18, refs, 8, C.MUTED)
    add_text(slide, 12.1, 6.88, 0.45, 0.18, f"{slide_no:02d}", 10, C.MUTED, True)


def add_panel(slide, x: float, y: float, w: float, h: float, fill: RGBColor = C.PANEL, border: RGBColor = C.BORDER):
    shp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(1)
    return shp


def add_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    size: int = 14,
    color: RGBColor = C.INK,
    bold: bool = False,
    font: str = "Aptos",
    align=PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = font


def add_bullets(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    items: Sequence[str],
    size: int = 14,
    color: RGBColor = C.INK,
    bullet_color: RGBColor = C.CYAN,
) -> None:
    rail = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.05), Inches(h)
    )
    rail.fill.solid()
    rail.fill.fore_color.rgb = bullet_color
    rail.line.fill.background()

    box = slide.shapes.add_textbox(Inches(x + 0.14), Inches(y), Inches(w - 0.14), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.bullet = True
        p.space_after = Pt(8)
        p.line_spacing = 1.08
        if p.runs:
            run = p.runs[0]
            run.font.name = "Aptos"
            run.font.size = Pt(size)
            run.font.color.rgb = color


def add_stat(slide, x: float, y: float, w: float, value: str, label: str, color: RGBColor) -> None:
    add_panel(slide, x, y, w, 0.84, C.PANEL, color)
    add_text(slide, x + 0.16, y + 0.08, w - 0.32, 0.24, value, 21, color, True)
    add_text(slide, x + 0.16, y + 0.43, w - 0.32, 0.18, label, 10, C.MUTED)


def add_chart(slide, x: float, y: float, w: float, h: float, categories: Sequence[str], values: Sequence[float],
              series_name: str, color: RGBColor) -> None:
    chart_data = CategoryChartData()
    chart_data.categories = list(categories)
    chart_data.add_series(series_name, list(values))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), chart_data
    ).chart
    chart.has_legend = False
    chart.value_axis.visible = True
    chart.value_axis.has_major_gridlines = True
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.format.line.color.rgb = C.BORDER
    chart.category_axis.format.line.color.rgb = C.BORDER
    chart.value_axis.major_gridlines.format.line.color.rgb = C.BORDER
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    plot.data_labels.font.size = Pt(10)
    plot.data_labels.font.bold = True
    ser = chart.series[0]
    fill = ser.format.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_table(slide, x: float, y: float, w: float, h: float, rows: Sequence[Sequence[str]], widths: Sequence[float]) -> None:
    table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h)).table
    total = sum(widths)
    for idx, frac in enumerate(widths):
        table.columns[idx].width = Inches(w * (frac / total))
    for r_idx, row in enumerate(rows):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = C.PANEL_ALT if r_idx == 0 else C.PANEL
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
            cell.text_frame.clear()
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = value
            r.font.name = "Aptos"
            r.font.size = Pt(10 if r_idx == 0 else 11)
            r.font.bold = r_idx == 0
            r.font.color.rgb = C.INK


def note(notes: list[str], slide_no: int, title: str, bullets: Iterable[str]) -> None:
    body = [f"Slide {slide_no:02d} - {title}"]
    body.extend(f"- {b}" for b in bullets)
    notes.append("\n".join(body))


def slide_title(prs: Presentation, notes: list[str]) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_tag(s, "EXECUTIVE THESIS", C.CYAN, C.BLUE_BG, 0.72, 0.5, 2.0)
    add_text(
        s,
        0.72,
        1.0,
        8.2,
        1.0,
        "Qualcomm has a credible shot to turn\nstrong edge AI silicon into a developer-preferred platform",
        28,
        C.INK,
        True,
        "Aptos Display",
    )
    add_text(
        s,
        0.72,
        2.05,
        8.6,
        0.65,
        "The winning move is not another set of isolated demos. It is a reusable, multimodal, OpenAI-style edge runtime on industrial-grade IQ-9075 hardware.",
        15,
        C.MUTED,
    )

    add_stat(s, 0.72, 3.15, 2.45, "4 public apps", "Public GenAI Studio modalities today", C.QCOM)
    add_stat(s, 3.42, 3.15, 2.45, "100 TOPS", "IQ-9075 public top configuration", C.GOLD)
    add_stat(s, 6.12, 3.15, 2.45, "$249 / 67 TOPS", "Jetson Orin Nano Super public entry point", C.JETSON)

    add_panel(s, 9.1, 1.0, 3.45, 4.3, C.PANEL, C.BORDER)
    add_text(s, 9.35, 1.3, 2.95, 0.22, "Truth model for the deck", 15, C.INK, True)
    add_bullets(
        s,
        9.35,
        1.75,
        2.95,
        2.95,
        [
            "Public: only what official Qualcomm and NVIDIA materials show now.",
            "Engineering evidence: current local repo capability, pending owner validation.",
            "Roadmap: next-step narrative, not public commitment.",
        ],
        13,
        C.INK,
        C.GOLD,
    )
    add_text(s, 9.35, 4.55, 2.95, 0.28, "Audience: exec, product, platform, GTM", 11, C.MUTED, True)
    add_footer(s, "[Q1][Q2][Q3][N1]", 1)
    note(notes, 1, "Executive thesis", [
        "Lead with platform story, not feature list.",
        "Establish the public/internal/roadmap truth model immediately.",
        "Frame Jetson as the benchmark for platform mindshare, not only hardware."
    ])


def slide_why_now(prs: Presentation, notes: list[str]) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        s,
        "Why this matters now",
        "Edge GenAI demand is moving from experimentation to deployment logic: local execution, privacy, low latency, and multimodal workflows.",
        PUBLIC,
        2,
        "[M1][M2][N6]",
    )
    add_panel(s, 0.72, 1.75, 5.25, 3.8, C.PANEL, C.BORDER)
    add_text(s, 0.95, 2.0, 4.7, 0.2, "Edge AI hardware market outlook", 15, C.INK, True)
    add_chart(s, 1.0, 2.35, 4.7, 2.65, ["2025", "2030"], [26.14, 58.90], "USD billions", C.QCOM)
    add_text(s, 1.15, 5.05, 4.3, 0.2, "MarketsandMarkets: edge AI hardware market, USD billions", 9, C.MUTED)

    add_panel(s, 6.25, 1.75, 6.35, 3.8, C.PANEL, C.BORDER)
    add_text(s, 6.48, 2.0, 5.8, 0.2, "What is driving urgency", 15, C.INK, True)
    add_table(
        s,
        6.45,
        2.35,
        5.9,
        2.9,
        [
            ["Driver", "Why it matters"],
            ["Privacy + control", "Sensitive voice, image, and operational data stays local."],
            ["Latency + resilience", "Field systems cannot depend on cloud round trips."],
            ["Cost discipline", "Local inference can reduce recurring API spend."],
            ["Multimodal workflows", "Real deployments combine text, voice, and vision."],
        ],
        [1.6, 4.3],
    )

    add_panel(s, 0.72, 5.88, 11.88, 0.46, C.DEEP, C.DEEP)
    add_text(
        s,
        0.92,
        5.98,
        11.4,
        0.18,
        "Implication: the next winning edge offer is not 'can it run a model?' but 'can teams ship multiple AI workflows on one local platform developers already know how to consume?'",
        12,
        RGBColor(255, 255, 255),
        True,
    )
    note(notes, 2, "Why now", [
        "Use one clean market chart, then pivot to concrete adoption drivers.",
        "Keep the slide focused on deployment logic rather than AI hype."
    ])


def slide_public_proof(prs: Presentation, notes: list[str]) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        s,
        "What Qualcomm publicly proves today",
        "Public material already shows a real containerized edge AI package. The gap is platform packaging, not lack of visible capability.",
        PUBLIC,
        3,
        "[Q1][Q2]",
    )
    add_panel(s, 0.72, 1.75, 7.0, 4.8, C.PANEL, C.BORDER)
    add_text(s, 0.95, 2.0, 6.5, 0.2, "Public proof points from GenAI Studio", 15, C.INK, True)
    add_table(
        s,
        0.95,
        2.3,
        6.5,
        3.65,
        [
            ["Public item", "Evidence visible today"],
            ["Containerized solution", "Official docs describe GenAI Studio as a containerized solution for Qualcomm Linux systems."],
            ["Supported public apps", "Text generation, text-to-speech, image generation, speech-to-text."],
            ["UI surface", "Official docs describe one web interface to access all apps."],
            ["Platforms", "IQ9 Qualcomm Linux and IQ9 Qualcomm Ubuntu are listed in public docs."],
            ["Developer workflow", "Public materials show model generation and container image build as part of the bring-up flow."],
        ],
        [2.0, 4.5],
    )

    add_panel(s, 8.0, 1.75, 4.6, 4.8, C.PANEL_ALT, C.BORDER)
    add_text(s, 8.24, 2.0, 4.1, 0.2, "Public architecture reading", 15, C.INK, True)
    layers = [
        ("Web UI", C.CYAN),
        ("Containerized app services", C.QCOM),
        ("Models + runtime prep", C.GOLD),
        ("IQ9 Linux / Ubuntu target", C.GREEN),
    ]
    y = 2.45
    for label, color in layers:
        add_panel(s, 8.32, y, 3.95, 0.62, C.PANEL, color)
        add_text(s, 8.54, y + 0.18, 3.4, 0.16, label, 14, C.INK, True)
        y += 0.82
    add_text(
        s,
        8.26,
        5.82,
        4.0,
        0.45,
        "Strategic read: Qualcomm already has public platform ingredients. It needs a stronger developer-platform narrative around them.",
        12,
        C.MUTED,
    )
    note(notes, 3, "Public proof", [
        "This slide proves Qualcomm has public credibility already.",
        "The key message is that the story exists but is not yet packaged as a platform story."
    ])


def slide_hardware(prs: Presentation, notes: list[str]) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        s,
        "Why IQ-9075 is a credible foundation",
        "Jetson wins the prototyping narrative. IQ-9075 gives Qualcomm a stronger industrial deployment narrative if the software experience is made simpler.",
        PUBLIC,
        4,
        "[Q3][Q4][Q5]",
    )
    add_stat(s, 0.72, 1.75, 2.45, "50 / 100 TOPS", "Public IQ-9075 SoC configurations", C.QCOM)
    add_stat(s, 3.4, 1.75, 2.45, "-40C to 115C", "Public thermal junction support", C.RED)
    add_stat(s, 6.08, 1.75, 2.45, "10+ years", "Public longevity support", C.GREEN)
    add_stat(s, 8.76, 1.75, 2.45, "Ubuntu + QLinux", "Public OS support statement", C.GOLD)

    add_panel(s, 0.72, 3.0, 12.0, 3.2, C.PANEL, C.BORDER)
    add_table(
        s,
        0.95,
        3.25,
        11.55,
        2.65,
        [
            ["Hardware signal", "Why it matters for platform strategy"],
            ["Industrial-grade positioning", "Supports a differentiated story in privacy-sensitive, always-on, field-deployed edge environments."],
            ["AI / robotics application emphasis", "Aligns naturally with voice, vision, and local copilot workflows near machines and operators."],
            ["EVK availability today", "Makes the platform story tangible for OEM, ODM, and developer evaluation."],
            ["Thermal + longevity support", "Strengthens Qualcomm's lane versus purely hobbyist or classroom-first positioning."],
        ],
        [2.7, 8.85],
    )
    note(notes, 4, "Hardware foundation", [
        "Tie every hardware fact back to software and deployment value.",
        "Avoid making this a generic silicon brag slide."
    ])


def slide_internal(prs: Presentation, notes: list[str]) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        s,
        "What the current repo already shows",
        "Use this only as internal or owner-validated material. It should not be implied as already public unless product and engineering owners confirm it.",
        INTERNAL,
        5,
        "[I1][I2][I3][I4][I5][I6]",
    )
    add_panel(s, 0.72, 1.75, 7.15, 4.65, C.PANEL, C.GOLD)
    add_text(s, 0.95, 2.0, 6.6, 0.2, "Current local platform evidence", 15, C.INK, True)
    add_table(
        s,
        0.95,
        2.3,
        6.6,
        3.55,
        [
            ["Module", "Current evidence in local repo"],
            ["Orchestrator + UI", "Canonical gateway on :8090 with unified status, helper routes, and OpenAI-style routing."],
            ["Text-Generation", "Direct service on :8088, routed via /v1/chat/completions and internal model mapping."],
            ["Speech-To-Text", "Direct service on :8081, routed via /v1/audio/transcriptions, /translations, stream, and /v1/realtime*."],
            ["Text-To-Speech", "Direct service on :8083, routed via /v1/audio/speech."],
            ["Image-Generation", "Backend on :8084, surfaced through /v1/images/generations with orchestrator arbitration."],
            ["Image-To-Text", "Vision flow via /v1/responses only, with in-container preprocessing."],
        ],
        [2.15, 4.45],
    )

    add_panel(s, 8.1, 1.75, 4.62, 4.65, C.PANEL_ALT, C.BORDER)
    add_text(s, 8.34, 2.0, 4.1, 0.2, "How to talk about this externally", 15, C.INK, True)
    add_bullets(
        s,
        8.34,
        2.38,
        4.02,
        2.05,
        [
            "Public today: containerized multimodal app package.",
            "Current internal milestone: unified OpenAI-style API surface across modalities.",
            "Validated prototype: multimodal routing and vision-assisted flows.",
            "Next release objective: one canonical external developer contract.",
        ],
        13,
        C.INK,
        C.GOLD,
    )
    add_panel(s, 8.34, 4.8, 3.9, 0.92, C.AMBER_BG, C.GOLD)
    add_text(
        s,
        8.54,
        5.03,
        3.5,
        0.4,
        "Accuracy rule: separate public evidence from engineering evidence on every customer-facing artifact.",
        11,
        C.INK,
        True,
    )
    note(notes, 5, "Internal capability", [
        "This slide is where the real platform opportunity becomes explicit.",
        "The right-hand column gives exact language discipline for executives and PMs."
    ])


def slide_aihub_growth(prs: Presentation, notes: list[str]) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        s,
        "Why Qualcomm AI Hub matters",
        "AI Hub is already a public software funnel for models, optimization, profiling, and sample apps. Linked well, it can feed GenAI Studio with easier model-to-solution workflows.",
        ROADMAP,
        6,
        "[H1][H2][H3][H4]",
    )
    add_stat(s, 0.72, 1.78, 2.65, "300+", "Public AI Hub model catalog", C.QCOM)
    add_stat(s, 3.58, 1.78, 2.65, "50+", "Hosted device types for evaluation", C.CYAN)
    add_stat(s, 6.44, 1.78, 2.65, "Funnel", "Public path from optimization to solution", C.GREEN)

    add_panel(s, 0.72, 2.95, 7.25, 3.2, C.PANEL, C.BORDER)
    add_text(s, 0.95, 3.15, 6.7, 0.2, "Public Qualcomm software funnel", 15, C.INK, True)
    stages = [
        ("AI Hub Models", "Discover optimized models", C.QCOM),
        ("Workbench", "Profile on hosted devices", C.CYAN),
        ("AI Hub Apps", "Show reusable app patterns", C.GOLD),
        ("GenAI Studio", "Deploy on IQ-9075", C.GREEN),
    ]
    x = 0.95
    for idx, (title, body, color) in enumerate(stages):
        add_panel(s, x, 3.55, 1.52, 1.72, C.PANEL_ALT if idx % 2 else C.PANEL, color)
        add_text(s, x + 0.1, 3.78, 1.28, 0.22, title, 11, C.INK, True)
        add_text(s, x + 0.1, 4.18, 1.28, 0.52, body, 9, C.MUTED)
        if idx < len(stages) - 1:
            chev = s.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(x + 1.56), Inches(4.14), Inches(0.34), Inches(0.34)
            )
            chev.fill.solid()
            chev.fill.fore_color.rgb = C.BORDER
            chev.line.fill.background()
        x += 1.86

    add_panel(s, 8.22, 2.95, 4.5, 3.2, C.PANEL_ALT, C.BORDER)
    add_text(s, 8.45, 3.15, 3.95, 0.2, "Strategic implication", 15, C.INK, True)
    add_bullets(
        s,
        8.45,
        3.48,
        3.85,
        2.0,
        [
            "Lower model-onboarding friction for teams evaluating Qualcomm first.",
            "Create a clearer story from model optimization to field deployment.",
            "Turn app demos into repeatable solution kits on Dragonwing hardware.",
            "Strengthen Qualcomm's software narrative without copying Jetson's exact playbook.",
        ],
        11,
        C.INK,
        C.CYAN,
    )
    add_panel(s, 8.45, 5.42, 3.85, 0.5, C.DEEP, C.DEEP)
    add_text(
        s,
        8.65,
        5.51,
        3.45,
        0.3,
        "Best read: AI Hub acquires developers; GenAI Studio helps them deploy on IQ-9075.",
        9,
        RGBColor(255, 255, 255),
        True,
    )
    note(notes, 6, "AI Hub growth lever", [
        "This slide connects Qualcomm's broader public software assets to the GenAI Studio deployment story.",
        "The message is model-to-device-to-solution, not just one more tool catalog."
    ])


def slide_unified_endpoints(prs: Presentation, notes: list[str]) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        s,
        "Why unifying many OpenAI-style endpoints matters",
        "Multiple local modalities become much more valuable when they are exposed through one familiar contract and one deployable runtime surface.",
        ROADMAP,
        7,
        "[Q1][Q2][I2][I6]",
    )
    boxes = [
        (0.88, "Hardware", "IQ-9075 / local edge deployment", C.QCOM),
        (3.0, "Runtime", "Models, containers, orchestration", C.CYAN),
        (5.12, "API", "Chat | audio | image | vision", C.GOLD),
        (7.24, "Apps", "Copilots, assistants, inspection, HMI", C.GREEN),
        (9.36, "Solutions", "Reusable customer workflows", C.PLUM),
    ]
    for x, title, body, color in boxes:
        add_panel(s, x, 2.3, 1.8, 2.0, C.PANEL, color)
        add_text(s, x + 0.15, 2.62, 1.5, 0.18, title, 16, C.INK, True)
        add_text(s, x + 0.15, 3.05, 1.5, 0.7, body, 11, C.MUTED)
        if x < 9.36:
            bar = s.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(x + 1.83), Inches(2.98), Inches(0.62), Inches(0.52)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = C.BORDER
            bar.line.fill.background()

    add_panel(s, 0.72, 5.0, 12.0, 1.25, C.PANEL_ALT, C.BORDER)
    add_table(
        s,
        0.95,
        5.18,
        11.55,
        0.86,
        [
            ["Developer payoff", "What changes"],
            ["Composition", "Chat, voice, image generation, and vision flows can be stitched without rebuilding integration patterns each time."],
            ["Migration", "Developers already familiar with OpenAI-style cloud APIs can port faster to edge deployment."],
        ],
        [1.7, 9.85],
    )
    note(notes, 7, "Unified endpoints", [
        "This is the strategic center of the whole deck.",
        "Make clear that endpoint unification is the difference between demos and platform primitives."
    ])


def slide_dev_value(prs: Presentation, notes: list[str]) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        s,
        "Developer value: move from demos to reusable platform primitives",
        "Jetson wins because its public software story is easy to understand. Qualcomm can close that gap by making GenAI Studio the canonical entry point.",
        ROADMAP,
        8,
        "[N2][N4][N5][I2]",
    )
    add_panel(s, 0.72, 1.8, 5.75, 4.8, C.PANEL, C.BORDER)
    add_text(s, 0.95, 2.02, 5.2, 0.2, "If GenAI Studio is treated like isolated demos", 15, C.RED, True)
    add_bullets(
        s,
        0.95,
        2.42,
        5.15,
        2.3,
        [
            "Each modality feels like a separate sample.",
            "Integration work repeats across customer projects.",
            "Cloud-to-edge migration remains a custom exercise.",
            "The platform story stays weaker than the silicon story.",
        ],
        14,
        C.INK,
        C.RED,
    )

    add_panel(s, 6.86, 1.8, 5.86, 4.8, C.PANEL, C.BORDER)
    add_text(s, 7.08, 2.02, 5.3, 0.2, "If GenAI Studio is treated like a platform", 15, C.GREEN, True)
    add_bullets(
        s,
        7.08,
        2.42,
        5.25,
        2.3,
        [
            "One API pattern across multiple local modalities.",
            "Reusable building blocks for partners and app teams.",
            "Faster prototyping and lower app-layer glue code.",
            "A cleaner public answer to 'why Qualcomm over Jetson?'.",
        ],
        14,
        C.INK,
        C.GREEN,
    )

    add_panel(s, 0.72, 5.92, 12.0, 0.44, C.DEEP, C.DEEP)
    add_text(
        s,
        0.92,
        6.01,
        11.5,
        0.16,
        "Platform unlock = one familiar surface for many edge workloads, backed by real hardware and deployment policy.",
        12,
        RGBColor(255, 255, 255),
        True,
    )
    note(notes, 8, "Developer value", [
        "This slide reframes the entire product as a developer experience decision.",
        "Keep the contrast simple and memorable."
    ])


def slide_customer_value(prs: Presentation, notes: list[str]) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        s,
        "Customer value and top use-case families",
        "The strongest story is not one killer demo. It is one local stack that can support many privacy-sensitive, low-latency, multimodal workflows.",
        ROADMAP,
        9,
        "[Q3][Q4][M2]",
    )
    add_panel(s, 0.72, 1.82, 7.15, 4.65, C.PANEL, C.BORDER)
    add_text(s, 0.95, 2.04, 6.6, 0.2, "Use-case matrix", 15, C.INK, True)
    add_table(
        s,
        0.95,
        2.35,
        6.6,
        3.75,
        [
            ["Use-case family", "Modalities", "Why local matters"],
            ["Operator copilot", "Chat + STT + TTS", "Hands-free, fast, private"],
            ["Inspection assistant", "Vision + chat + TTS", "Image understanding near equipment"],
            ["Secure transcription", "STT + chat", "Local data control, lower cloud cost"],
            ["Field workflow assistant", "Vision + voice + text", "Offline resilience and low latency"],
            ["Creative local asset flow", "Image generation + chat", "Rapid local iteration"],
        ],
        [2.1, 1.75, 2.75],
    )

    add_panel(s, 8.15, 1.82, 4.57, 4.65, C.PANEL_ALT, C.BORDER)
    add_text(s, 8.38, 2.04, 4.05, 0.2, "What customers actually buy", 15, C.INK, True)
    add_bullets(
        s,
        8.38,
        2.44,
        3.95,
        2.35,
        [
            "Low latency and predictable local UX.",
            "Privacy, data residency, and local control.",
            "Reduced recurring inference spend.",
            "Simpler integration for teams already used to cloud API patterns.",
        ],
        14,
        C.INK,
        C.CYAN,
    )
    add_panel(s, 8.38, 5.18, 3.95, 0.86, C.BLUE_BG, C.CYAN)
    add_text(
        s,
        8.58,
        5.42,
        3.55,
        0.3,
        "Messaging discipline: say 'customer demand signals' unless direct customer evidence is formally validated.",
        11,
        C.INK,
        True,
    )
    note(notes, 9, "Customer value", [
        "Anchor the story in local latency, privacy, resilience, and reuse.",
        "Use industrially relevant examples, not generic chatbot examples."
    ])


def slide_jetson_advantage(prs: Presentation, notes: list[str]) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        s,
        "Why Jetson wins mindshare today",
        "Jetson's edge advantage is not only hardware. It is the visible software funnel: entry price, JetPack, tutorials, Isaac ROS, forums, and AI Lab content.",
        PUBLIC,
        10,
        "[N1][N2][N4][N5][N6][N7]",
    )
    add_stat(s, 0.72, 1.78, 2.45, "$249", "Jetson Orin Nano Super public dev-kit price", C.JETSON)
    add_stat(s, 3.42, 1.78, 2.45, "67 TOPS", "Jetson Orin Nano Super public AI performance", C.JETSON)
    add_stat(s, 6.12, 1.78, 2.45, "JetPack", "Official software stack brand", C.JETSON)
    add_stat(s, 8.82, 1.78, 2.45, "Isaac ROS", "Visible robotics ecosystem anchor", C.JETSON)

    stages = [
        ("Low-cost entry", "$249 dev kit lowers barrier"),
        ("Software stack", "JetPack, CUDA, TensorRT, cuDNN"),
        ("Learning funnel", "Getting started, DLI, tutorials, AI Lab"),
        ("Robotics stack", "Isaac ROS and related workflows"),
        ("Community pull", "Forums, projects, ecosystem momentum"),
    ]
    x = 0.95
    for idx, (title, body) in enumerate(stages):
        add_panel(s, x, 3.05, 2.15, 2.1, C.PANEL if idx % 2 == 0 else C.PANEL_ALT, C.JETSON)
        add_text(s, x + 0.14, 3.28, 1.85, 0.25, title, 13, C.INK, True)
        add_text(s, x + 0.14, 3.72, 1.85, 0.7, body, 11, C.MUTED)
        if idx < len(stages) - 1:
            chev = s.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(x + 2.2), Inches(3.83), Inches(0.42), Inches(0.45)
            )
            chev.fill.solid()
            chev.fill.fore_color.rgb = C.JETSON
            chev.line.fill.background()
        x += 2.42

    add_panel(s, 0.72, 5.55, 12.0, 0.52, C.DEEP, C.DEEP)
    add_text(
        s,
        0.92,
        5.68,
        11.5,
        0.16,
        "Jetson's public moat is software packaging and developer momentum. Qualcomm's response must be platform clarity, not only better silicon claims.",
        12,
        RGBColor(255, 255, 255),
        True,
    )
    note(notes, 10, "Jetson public advantage", [
        "This slide should feel fair and evidence-based.",
        "The point is public ecosystem gravity, not superiority in every workload."
    ])


def slide_qcom_gaps(prs: Presentation, notes: list[str]) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        s,
        "Where Qualcomm is lagging publicly",
        "This is a public positioning gap analysis. It does not imply Qualcomm lacks the technical foundation to close the gap.",
        ROADMAP,
        11,
        "[Q1][Q2][Q3][N2][N4][N5]",
    )
    add_panel(s, 0.72, 1.82, 12.0, 4.85, C.PANEL, C.BORDER)
    add_table(
        s,
        0.95,
        2.05,
        11.55,
        4.15,
        [
            ["Public gap", "Why it hurts adoption", "What would fix it"],
            ["Lower developer mindshare", "Fewer teams start with Qualcomm by default", "Make GenAI Studio the obvious first touchpoint"],
            ["Less visible onboarding simplicity", "More setup friction in early evaluation", "Publish one canonical quickstart and validation path"],
            ["App-bundle perception", "Customers see demos more than platform primitives", "Lead with reusable APIs and reference solutions"],
            ["Weaker software-platform branding", "Great silicon can be undervalued", "Create a consistent runtime and developer platform narrative"],
            ["Lower public community momentum", "Ecosystem compounding works against Qualcomm", "Invest in tutorials, partner kits, and migration stories"],
        ],
        [2.7, 4.05, 4.8],
    )
    note(notes, 11, "Qualcomm public gaps", [
        "Keep the tone constructive. This is fixable positioning and productization work.",
        "Do not phrase this as a hardware weakness slide."
    ])


def slide_gap_closure(prs: Presentation, notes: list[str]) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        s,
        "How this solution reduces or eliminates those gaps",
        "The core move is to convert capability into platform: familiar API surface, multimodal composition, and an easier path from prototype to deployment.",
        ROADMAP,
        12,
        "[I2][I6][Q3][H1][H4][N2]",
    )
    add_table(
        s,
        0.72,
        1.85,
        12.0,
        4.2,
        [
            ["Today", "Platform response", "Effect"],
            ["Separate public apps", "One orchestrated multimodal entry point", "Stronger product coherence"],
            ["Higher integration friction", "One OpenAI-style contract across workloads", "Faster developer adoption"],
            ["Weak reuse story", "Composable edge primitives for voice, text, image, vision", "More customer use cases on one stack"],
            ["Hardware-led message", "Hardware -> runtime -> API -> solution narrative", "Better executive and GTM story"],
            ["Jetson mindshare lead", "Sharper platform positioning on industrial local AI", "More credible competitive answer"],
        ],
        [2.6, 5.0, 4.4],
    )
    add_panel(s, 0.72, 6.25, 12.0, 0.42, C.DEEP, C.DEEP)
    add_text(
        s,
        0.92,
        6.33,
        11.5,
        0.16,
        "Do not position this as 'Jetson, but Qualcomm.' Position it as the easier software layer for rugged, privacy-sensitive, multimodal edge AI on industrial hardware.",
        12,
        RGBColor(255, 255, 255),
        True,
    )
    note(notes, 12, "Gap closure", [
        "This is the answer slide.",
        "Translate every weakness into a platform action and a business effect."
    ])


def slide_competitive_snapshot(prs: Presentation, notes: list[str]) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        s,
        "Competitive snapshot: Qualcomm vs Jetson",
        "Jetson leads in public ecosystem maturity. Qualcomm can win on industrial-grade local multimodal platformization if it makes the software story easier to adopt.",
        ROADMAP,
        13,
        "[Q3][Q4][N1][N2][N7]",
    )
    add_table(
        s,
        0.72,
        1.8,
        12.0,
        4.8,
        [
            ["Dimension", "Qualcomm GenAI Studio on IQ-9075", "Jetson Orin Nano / Jetson stack"],
            ["Public identity", "Industrial edge AI platform emerging", "Developer-first edge AI ecosystem established"],
            ["Entry on-ramp", "Needs simpler canonical journey", "Very strong price + tutorial + JetPack funnel"],
            ["Industrial deployment fit", "Strong public thermal, longevity, OS, and rugged deployment story", "Strong prototyping story; industrial fit often comes through modules/partners"],
            ["Unified multimodal API story", "High potential; current internal direction appears strong", "Broad model support and tutorials; public story centered on JetPack/CUDA ecosystem"],
            ["Robotics ecosystem", "Less visible publicly", "Strong public position via Isaac ROS and related workflows"],
            ["Best strategic lane", "Privacy-sensitive, rugged, reusable local AI platform", "Broad developer mindshare and ecosystem breadth"],
        ],
        [2.1, 4.8, 5.1],
    )
    note(notes, 13, "Competitive snapshot", [
        "This is a fair comparison slide, not a hype slide.",
        "The key is to show a credible Qualcomm lane without pretending Jetson's ecosystem lead does not exist."
    ])


def slide_priorities(prs: Presentation, notes: list[str]) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        s,
        "Recommended next actions",
        "The fastest path to closing the platform gap is to simplify the public software story and make the API surface feel canonical.",
        ROADMAP,
        14,
        "[Q1][Q2][H1][H3][I2][I6]",
    )
    cards = [
        ("1. Reposition publicly", "Shift GenAI Studio messaging from 'four apps' to 'multimodal edge AI platform starter kit'.", C.QCOM),
        ("2. Publish one contract", "Stabilize and document one canonical external API surface with clear GA vs preview labels.", C.CYAN),
        ("3. Simplify onboarding", "Ship one happy-path Ubuntu IQ-9075 quickstart and one validation command.", C.GOLD),
        ("4. Prove reuse", "Publish 2-3 reference solutions built from the same primitives and connect them to AI Hub assets.", C.GREEN),
        ("5. Equip field teams", "Create a Jetson response kit, migration notes, and partner enablement story.", C.PLUM),
    ]
    y = 1.85
    for title, body, color in cards:
        add_panel(s, 0.88, y, 11.6, 0.73, C.PANEL, color)
        add_text(s, 1.1, y + 0.14, 2.3, 0.18, title, 14, C.INK, True)
        add_text(s, 3.35, y + 0.14, 8.75, 0.18, body, 12, C.MUTED)
        y += 0.86

    add_panel(s, 0.88, 6.18, 11.6, 0.42, C.DEEP, C.DEEP)
    add_text(
        s,
        1.08,
        6.26,
        11.2,
        0.16,
        "Executive ask: align product, engineering, docs, and GTM around one validated public platform narrative and one canonical developer journey.",
        12,
        RGBColor(255, 255, 255),
        True,
    )
    note(notes, 14, "Recommended actions", [
        "Keep the asks platform-oriented and executable.",
        "Leadership should be able to assign owners directly from this slide."
    ])


def slide_appendix(prs: Presentation, notes: list[str]) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        s,
        "Appendix: validation boundaries and source map",
        "Use this slide to preserve claim discipline whenever the deck is reused.",
        PUBLIC,
        15,
        "[Q1][Q2][Q3][H1][H3][N1][N2][M1][M2][I1][I2][I6]",
    )
    add_panel(s, 0.72, 1.82, 3.85, 4.8, C.PANEL, C.BORDER)
    add_text(s, 0.95, 2.04, 3.35, 0.2, "Public Qualcomm sources", 15, C.INK, True)
    add_text(
        s,
        0.95,
        2.34,
        3.3,
        3.9,
        "Q1 GenAI Studio docs\nQ2 Public GenAI Studio GitHub README\nQ3 IQ-9075 product page\nQ4 IQ-9075 EVK page\nQ5 IQ-9075 EVK brief PDF\nH1 AI Hub home\nH3 AI Hub Apps",
        12,
        C.INK,
    )

    add_panel(s, 4.76, 1.82, 3.95, 4.8, C.PANEL, C.BORDER)
    add_text(s, 4.99, 2.04, 3.45, 0.2, "Public NVIDIA + market sources", 15, C.INK, True)
    add_text(
        s,
        4.99,
        2.34,
        3.4,
        4.1,
        "N1 Jetson Orin Nano Super page\nN2 JetPack official page\nN4 Jetson getting started\nN5 Tutorials / GenAI Lab\nN6 Jetson LLM/VLM blog\nN7 Isaac ROS\nM1 MarketsandMarkets\nM2 Global Market Insights\nM3 SHD report release article",
        12,
        C.INK,
    )

    add_panel(s, 8.92, 1.82, 3.8, 4.8, C.PANEL_ALT, C.BORDER)
    add_text(s, 9.15, 2.04, 3.3, 0.2, "Validation rules", 15, C.INK, True)
    add_bullets(
        s,
        9.15,
        2.34,
        3.2,
        3.9,
        [
            "Only public Qualcomm sources should support public-available claims.",
            "Local repo evidence should be labeled engineering evidence until owner-validated.",
            "Roadmap statements should be framed as next-step objectives, not public commitments.",
            "Avoid unsupported market-share claims unless explicitly sourced.",
        ],
        12,
        C.INK,
        C.GOLD,
    )
    note(notes, 15, "Appendix", [
        "This slide is the accuracy backstop for the whole deck.",
        "It protects the narrative from drifting into unsupported public claims."
    ])


if __name__ == "__main__":
    build()
    print(f"Wrote {OUT_PPTX}")
    print(f"Wrote {OUT_NOTES}")
    print(f"Wrote {OUT_SOURCES}")
