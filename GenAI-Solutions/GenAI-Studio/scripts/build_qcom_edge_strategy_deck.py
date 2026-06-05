#!/usr/bin/env python3
"""
Build an executive PowerPoint positioning Qualcomm GenAI Studio on IQ-9075 EVK
versus NVIDIA Jetson Orin Nano / Jetson edge platforms.

This deck intentionally separates:
1) Public proof visible today
2) Current engineering evidence from the local repo
3) Roadmap / next-step narrative
"""

from __future__ import annotations

from pathlib import Path
from typing import Iterable, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "artifacts"
OUT_DIR.mkdir(exist_ok=True)
OUT_PPTX = OUT_DIR / "Qualcomm_GenAI_Studio_IQ9075_vs_Jetson_Strategy_2026-04-28.pptx"
OUT_NOTES = OUT_DIR / "Qualcomm_GenAI_Studio_IQ9075_vs_Jetson_Strategy_2026-04-28_notes.md"


class C:
    BG = RGBColor(10, 16, 32)
    BG_2 = RGBColor(16, 25, 46)
    PANEL = RGBColor(18, 28, 52)
    PANEL_2 = RGBColor(24, 38, 68)
    WHITE = RGBColor(245, 248, 255)
    MUTED = RGBColor(177, 190, 214)
    BORDER = RGBColor(55, 71, 110)
    CYAN = RGBColor(98, 210, 255)
    TEAL = RGBColor(84, 219, 191)
    GOLD = RGBColor(246, 196, 83)
    CORAL = RGBColor(255, 134, 120)
    LIME = RGBColor(118, 185, 0)
    PURPLE = RGBColor(177, 124, 255)


PUBLIC_TAG = ("PUBLIC", C.CYAN)
INTERNAL_TAG = ("ENGINEERING EVIDENCE", C.GOLD)
ROADMAP_TAG = ("ROADMAP", C.PURPLE)
INFERENCE_TAG = ("INFERENCE", C.CORAL)

SLIDE_W = 13.333
SLIDE_H = 7.5


def build_deck() -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    notes: list[str] = []

    slide_1_title(prs, notes)
    slide_2_why_now(prs, notes)
    slide_3_public_proof(prs, notes)
    slide_4_hardware_foundation(prs, notes)
    slide_5_internal_capability(prs, notes)
    slide_6_unified_api_value(prs, notes)
    slide_7_developer_value(prs, notes)
    slide_8_customer_value(prs, notes)
    slide_9_jetson_mindshare(prs, notes)
    slide_10_qualcomm_public_gaps(prs, notes)
    slide_11_gap_closure(prs, notes)
    slide_12_compare_matrix(prs, notes)
    slide_13_gtm(prs, notes)
    slide_14_priorities(prs, notes)
    slide_15_appendix(prs, notes)

    prs.save(str(OUT_PPTX))
    OUT_NOTES.write_text("\n\n".join(notes), encoding="utf-8")


def add_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = C.BG
    top = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_W), Inches(0.34)
    )
    top.fill.solid()
    top.fill.fore_color.rgb = C.BG_2
    top.line.fill.background()


def add_title(slide, title: str, subtitle: str = "", tag: tuple[str, RGBColor] | None = None) -> None:
    add_bg(slide)
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.52), Inches(8.7), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = C.WHITE
    r.font.name = "Aptos Display"
    if subtitle:
        sp = tf.add_paragraph()
        sr = sp.add_run()
        sr.text = subtitle
        sr.font.size = Pt(14)
        sr.font.color.rgb = C.MUTED
        sr.font.name = "Aptos"
    if tag is not None:
        add_tag(slide, tag[0], tag[1], 10.75, 0.52, 2.0, 0.42)


def add_footer(slide, sources: str, slide_no: int) -> None:
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.45), Inches(6.84), Inches(12.45), Inches(0.01)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = C.BORDER
    line.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.55), Inches(6.88), Inches(11.1), Inches(0.35))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = sources
    r.font.size = Pt(8)
    r.font.color.rgb = C.MUTED
    r.font.name = "Aptos"

    num = slide.shapes.add_textbox(Inches(12.0), Inches(6.85), Inches(0.75), Inches(0.3))
    ntf = num.text_frame
    ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
    np = ntf.paragraphs[0]
    np.alignment = PP_ALIGN.RIGHT
    nr = np.add_run()
    nr.text = f"{slide_no:02d}"
    nr.font.size = Pt(11)
    nr.font.bold = True
    nr.font.color.rgb = C.MUTED
    nr.font.name = "Aptos"


def add_tag(slide, text: str, color: RGBColor, x: float, y: float, w: float, h: float) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = C.BG
    r.font.name = "Aptos"


def add_panel(slide, x: float, y: float, w: float, h: float, title: str = "", fill: RGBColor = C.PANEL):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = C.BORDER
    shape.line.width = Pt(1)
    if title:
        tf = shape.text_frame
        tf.clear()
        tf.margin_left = Pt(10)
        tf.margin_right = Pt(10)
        tf.margin_top = Pt(8)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = C.WHITE
        r.font.name = "Aptos"
    return shape


def add_textbox(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    font_size: int = 16,
    color: RGBColor = C.WHITE,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    name: str = "Aptos",
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = name


def add_bullets(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    bullets: Sequence[str],
    font_size: int = 17,
    color: RGBColor = C.WHITE,
    bullet_color: RGBColor = C.CYAN,
    spacing: float = 1.12,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, item in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.line_spacing = spacing
        p.space_after = Pt(8)
        p.bullet = True
        if p.runs:
            run = p.runs[0]
        else:
            run = p.add_run()
            run.text = item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Aptos"
    # python-pptx does not expose bullet color directly; add accent rail beside the box
    rail = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x - 0.08), Inches(y + 0.02), Inches(0.02), Inches(h - 0.04)
    )
    rail.fill.solid()
    rail.fill.fore_color.rgb = bullet_color
    rail.line.fill.background()


def add_stat_card(slide, x: float, y: float, w: float, h: float, value: str, label: str, color: RGBColor) -> None:
    add_panel(slide, x, y, w, h, fill=C.PANEL_2)
    add_textbox(slide, x + 0.15, y + 0.12, w - 0.3, 0.45, value, font_size=24, color=color, bold=True)
    add_textbox(slide, x + 0.15, y + 0.58, w - 0.3, 0.34, label, font_size=11, color=C.MUTED)


def add_two_col_compare(slide, x: float, y: float, w: float, h: float, left_title: str, right_title: str,
                        left_items: Sequence[str], right_items: Sequence[str],
                        left_color: RGBColor, right_color: RGBColor) -> None:
    gap = 0.2
    col_w = (w - gap) / 2
    add_panel(slide, x, y, col_w, h, title=left_title, fill=C.PANEL)
    add_panel(slide, x + col_w + gap, y, col_w, h, title=right_title, fill=C.PANEL)
    add_bullets(slide, x + 0.18, y + 0.5, col_w - 0.3, h - 0.65, left_items, font_size=15, bullet_color=left_color)
    add_bullets(
        slide, x + col_w + gap + 0.18, y + 0.5, col_w - 0.3, h - 0.65, right_items, font_size=15, bullet_color=right_color
    )


def add_matrix_table(slide, x: float, y: float, w: float, h: float, rows: Sequence[Sequence[str]],
                     col_widths: Sequence[float], header_fill: RGBColor = C.PANEL_2) -> None:
    table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h)).table
    total = sum(col_widths)
    for idx, frac in enumerate(col_widths):
        table.columns[idx].width = Inches(w * (frac / total))
    row_h = h / len(rows)
    for row in table.rows:
        row.height = Inches(row_h)
    for r_idx, row_vals in enumerate(rows):
        for c_idx, value in enumerate(row_vals):
            cell = table.cell(r_idx, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill if r_idx == 0 else C.PANEL
            cell.text_frame.clear()
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = value
            run.font.name = "Aptos"
            run.font.size = Pt(11 if r_idx == 0 else 12)
            run.font.bold = r_idx == 0
            run.font.color.rgb = C.WHITE if r_idx == 0 else C.WHITE
            cell.margin_left = Pt(6)
            cell.margin_right = Pt(6)
            cell.margin_top = Pt(4)
            cell.margin_bottom = Pt(4)


def add_arrow(slide, x1: float, y1: float, x2: float, y2: float, color: RGBColor = C.CYAN) -> None:
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = color
    line.line.width = Pt(2.2)
    line.line.end_arrowhead = True


def add_notes(notes: list[str], slide_no: int, title: str, bullets: Iterable[str]) -> None:
    lines = [f"Slide {slide_no:02d} - {title}"]
    lines.extend(f"- {item}" for item in bullets)
    notes.append("\n".join(lines))


def slide_1_title(prs: Presentation, notes: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    hero = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(0.68), Inches(12.15), Inches(5.75)
    )
    hero.fill.solid()
    hero.fill.fore_color.rgb = C.PANEL
    hero.line.color.rgb = C.BORDER
    hero.line.width = Pt(1.2)

    add_tag(slide, "EXECUTIVE THESIS", C.CYAN, 0.82, 0.94, 1.9, 0.42)
    add_textbox(
        slide,
        0.82,
        1.42,
        7.5,
        1.0,
        "From strong edge AI silicon to\n"
        "developer-preferred multimodal edge platform",
        font_size=28,
        color=C.WHITE,
        bold=True,
        name="Aptos Display",
    )
    add_textbox(
        slide,
        0.84,
        2.52,
        7.0,
        0.8,
        "Qualcomm GenAI Studio on Dragonwing IQ-9075 EVK can become Qualcomm's public answer to Jetson's platform advantage "
        "if it is positioned and productized as a reusable OpenAI-compatible edge runtime, not just a demo bundle.",
        font_size=15,
        color=C.MUTED,
    )

    add_stat_card(slide, 0.86, 3.6, 2.35, 0.95, "4 public apps", "Text, speech, image generation, transcription in public GenAI Studio docs", C.CYAN)
    add_stat_card(slide, 3.42, 3.6, 2.35, 0.95, "100 TOPS", "IQ-9075 top configuration; industrial-grade positioning", C.TEAL)
    add_stat_card(slide, 5.98, 3.6, 2.35, 0.95, "1 API style", "Internal strategy opportunity: familiar OpenAI-compatible surface", C.GOLD)

    add_panel(slide, 8.72, 1.08, 3.1, 3.9, title="Truth layers", fill=C.PANEL_2)
    add_bullets(
        slide,
        8.94,
        1.58,
        2.7,
        2.8,
        [
            "Public proof today: only what official Qualcomm repo/docs show now.",
            "Engineering evidence: local repo capabilities, pending owner validation.",
            "Roadmap narrative: next-release platform story, not public commitment.",
        ],
        font_size=14,
        bullet_color=C.GOLD,
    )

    add_panel(slide, 8.72, 5.22, 3.1, 0.82, fill=C.BG_2)
    add_textbox(slide, 8.94, 5.42, 2.7, 0.28, "Audience: executives, product, platform, GTM", font_size=12, color=C.WHITE, bold=True)
    add_textbox(slide, 8.94, 5.68, 2.7, 0.2, "Date basis: April 28, 2026", font_size=10, color=C.MUTED)

    add_footer(
        slide,
        "Public sources: Qualcomm GenAI Studio docs/repo, Qualcomm IQ-9075 pages, NVIDIA Jetson official pages. "
        "Internal references called out explicitly in later slides.",
        1,
    )
    add_notes(
        notes,
        1,
        "Title / executive thesis",
        [
            "Lead with platform thesis, not feature inventory.",
            "Frame this as a software platform unlock built on credible industrial hardware.",
            "Establish the three-layer truth model up front to avoid public/internal confusion.",
        ],
    )


def slide_2_why_now(prs: Presentation, notes: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        slide,
        "Why now: edge GenAI is shifting from experiments to deployment logic",
        "The market pull is no longer just faster models. It is local execution, privacy, resiliency, and multimodal workflows at the edge.",
        PUBLIC_TAG,
    )

    add_stat_card(slide, 0.62, 1.42, 2.3, 0.95, "24.8% CAGR", "Edge AI market growth, 2025-2034 (Global Market Insights)", C.CYAN)
    add_stat_card(slide, 3.08, 1.42, 2.45, 0.95, "$58.9B by 2030", "Edge AI hardware market projection (MarketsandMarkets)", C.TEAL)
    add_stat_card(slide, 5.69, 1.42, 2.45, 0.95, "Local first", "Latency, privacy, and resilience are becoming design requirements", C.GOLD)

    force_titles = [
        "Latency and resilience",
        "Privacy and data control",
        "Cloud cost pressure",
        "Multimodal field workflows",
    ]
    force_text = [
        "Industrial and robotics workflows cannot always depend on cloud round trips or persistent connectivity.",
        "Local processing reduces movement of sensitive audio, image, and operational data off device.",
        "On-device inference can cap recurring API spend and make edge deployments economically reusable.",
        "Voice, text, and vision increasingly need to work together inside one operator flow, not as separate demos.",
    ]
    colors = [C.CYAN, C.TEAL, C.GOLD, C.CORAL]
    x_positions = [0.62, 3.75, 6.88, 10.01]
    for idx in range(4):
        add_panel(slide, x_positions[idx], 2.76, 2.66, 2.55, title=force_titles[idx], fill=C.PANEL)
        add_textbox(slide, x_positions[idx] + 0.14, 3.2, 2.38, 1.85, force_text[idx], font_size=13, color=C.WHITE)
        rail = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x_positions[idx]), Inches(2.76), Inches(0.08), Inches(2.55)
        )
        rail.fill.solid()
        rail.fill.fore_color.rgb = colors[idx]
        rail.line.fill.background()

    add_panel(slide, 0.62, 5.64, 12.0, 0.8, fill=C.BG_2)
    add_textbox(
        slide,
        0.82,
        5.87,
        11.6,
        0.28,
        "Implication: the winning edge AI offer is no longer 'can it run a model?' It is 'can teams ship multiple AI workflows on one local platform developers already know how to consume?'",
        font_size=13,
        color=C.WHITE,
        bold=True,
    )

    add_footer(
        slide,
        "Sources: Global Market Insights edge AI market (Mar 2025); MarketsandMarkets edge AI hardware press release; EDGE AI Foundation 2025 report.",
        2,
    )
    add_notes(
        notes,
        2,
        "Why now",
        [
            "Use one market metric and one hardware metric, then pivot quickly to concrete buying drivers.",
            "Avoid generic AI hype. Emphasize privacy, offline operation, latency, and multimodal workflows.",
            "This slide sets up why endpoint unification matters later.",
        ],
    )


def slide_3_public_proof(prs: Presentation, notes: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        slide,
        "What Qualcomm GenAI Studio publicly proves today",
        "Public documentation already shows a real containerized edge AI package. The public story is credible, but still framed more as a set of apps than as a developer platform.",
        PUBLIC_TAG,
    )

    add_two_col_compare(
        slide,
        0.62,
        1.45,
        7.0,
        4.95,
        "Public proof points",
        "What this means strategically",
        [
            "Containerized GenAI Studio for Qualcomm Linux systems.",
            "Four publicly documented use cases: text generation, text-to-speech, image generation, speech-to-text.",
            "Web interface documented as a unified user entry point.",
            "Published service ports: asr 8081, text-to-image 8084, text-to-text 8088, text-to-speech 8083.",
            "Publicly documented models: Llama 3.2-3B, Melo-TTS, Stable Diffusion, Whisper Base.",
        ],
        [
            "Qualcomm is already demonstrating a multi-container software packaging model, not a single binary demo.",
            "The public materials already imply workflow integration and practical prototyping intent.",
            "The current public gap is not 'no product' but 'insufficient platform-level narrative'.",
            "This is enough foundation to tell a stronger developer-platform story if API and onboarding are unified.",
        ],
        C.CYAN,
        C.TEAL,
    )

    add_panel(slide, 7.88, 1.45, 4.72, 4.95, title="Public architecture shape", fill=C.PANEL)
    layers = [
        ("Web UI", C.CYAN),
        ("Containerized app services", C.TEAL),
        ("Models + runtime + device setup", C.GOLD),
        ("IQ9 Qualcomm Linux / Ubuntu target", C.CORAL),
    ]
    y = 2.0
    for text, color in layers:
        add_panel(slide, 8.3, y, 3.86, 0.7, fill=C.PANEL_2)
        add_textbox(slide, 8.55, y + 0.19, 3.35, 0.22, text, font_size=15, color=C.WHITE, bold=True)
        rail = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(8.3), Inches(y), Inches(0.09), Inches(0.7)
        )
        rail.fill.solid()
        rail.fill.fore_color.rgb = color
        rail.line.fill.background()
        y += 0.92
    add_arrow(slide, 10.22, 2.72, 10.22, 2.92)
    add_arrow(slide, 10.22, 3.64, 10.22, 3.84)
    add_arrow(slide, 10.22, 4.56, 10.22, 4.76)

    add_footer(
        slide,
        "Sources: Qualcomm GenAI Studio doc (updated Feb 23, 2026); public GenAI Studio GitHub README.",
        3,
    )
    add_notes(
        notes,
        3,
        "Public proof today",
        [
            "Make the distinction: public evidence already proves platform intent, but not yet full platform branding.",
            "Public repo/docs are enough to show containerization, four modalities, web UI, ports, and supported models.",
            "This slide is the bridge from credibility to opportunity.",
        ],
    )


def slide_4_hardware_foundation(prs: Presentation, notes: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        slide,
        "Why IQ-9075 matters: industrial-grade hardware makes the platform story different",
        "Jetson often wins the prototyping conversation. IQ-9075 gives Qualcomm a stronger story where customers care about ruggedness, longevity, local AI, and field deployment.",
        PUBLIC_TAG,
    )

    add_stat_card(slide, 0.62, 1.5, 2.2, 0.95, "50 / 100 TOPS", "IQ-9075 SoC configurations", C.CYAN)
    add_stat_card(slide, 2.98, 1.5, 2.35, 0.95, "-40C to 115C", "Thermal junction support in IQ9 public product page", C.TEAL)
    add_stat_card(slide, 5.49, 1.5, 2.2, 0.95, "10+ years", "Product longevity support", C.GOLD)
    add_stat_card(slide, 7.85, 1.5, 2.1, 0.95, "Ubuntu + Qualcomm Linux", "Publicly documented software support", C.CORAL)

    add_panel(slide, 0.62, 2.85, 5.75, 2.95, title="What Qualcomm hardware pages emphasize", fill=C.PANEL)
    add_bullets(
        slide,
        0.82,
        3.32,
        5.35,
        2.2,
        [
            "Demanding industrial applications and extreme environments.",
            "AI/ML and robotics workflows on the IQ-9075 EVK.",
            "Development hardware available now for evaluation and prototyping.",
            "High-compute, power-efficient AI positioning rather than hobbyist-first positioning.",
        ],
        font_size=15,
        bullet_color=C.CYAN,
    )

    add_panel(slide, 6.65, 2.85, 5.95, 2.95, title="Strategic implication", fill=C.PANEL_2)
    add_textbox(
        slide,
        6.9,
        3.34,
        5.45,
        0.96,
        "Qualcomm does not need to beat Jetson by looking more like a student dev kit.\n"
        "It needs to win by turning industrial-grade local AI hardware into an easier, more familiar software platform.",
        font_size=17,
        color=C.WHITE,
        bold=True,
    )
    add_bullets(
        slide,
        6.9,
        4.45,
        5.35,
        1.15,
        [
            "Better fit for privacy-sensitive enterprise and industrial deployments.",
            "Natural foundation for rugged multimodal edge AI, not only maker demos.",
        ],
        font_size=14,
        bullet_color=C.GOLD,
    )

    add_footer(
        slide,
        "Sources: Qualcomm IQ-9075 product page; IQ-9075 EVK developer page; IQ9 documentation and EVK quickstart pages.",
        4,
    )
    add_notes(
        notes,
        4,
        "Hardware foundation",
        [
            "This slide explains why Qualcomm can differentiate without copying Jetson's public identity.",
            "Anchor the argument in industrial readiness, OS support, longevity, and robotics/AI workflows.",
            "Keep the hardware message tied to software/platform value.",
        ],
    )


def slide_5_internal_capability(prs: Presentation, notes: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        slide,
        "What current engineering capability appears to add",
        "This slide is not a public claim. It is based on current local engineering docs and code, and should be validated with product and engineering owners before any external use.",
        INTERNAL_TAG,
    )

    caution = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.62), Inches(1.35), Inches(12.0), Inches(0.62)
    )
    caution.fill.solid()
    caution.fill.fore_color.rgb = RGBColor(73, 52, 0)
    caution.line.color.rgb = C.GOLD
    tf = caution.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Validation boundary: local repo evidence only. Use labels such as 'current internal milestone' or 'validated prototype' until owners confirm public status."
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = C.WHITE
    r.font.name = "Aptos"

    add_two_col_compare(
        slide,
        0.62,
        2.2,
        6.25,
        4.25,
        "Internal platform evidence",
        "Why it matters",
        [
            "Orchestrator + UI on port 8090.",
            "Unified routes include /v1/models, /v1/chat/completions, /v1/audio/transcriptions, /v1/audio/translations, /v1/audio/speech, /v1/images/generations, and STT realtime routes.",
            "Image-to-Text routes exist for preprocess, vision, chat, and reset flows.",
            "OpenAI router can route vision-intent payloads into I2T.",
            "Current docs describe production traffic through orchestrator and NPU arbitration for mixed image workloads.",
        ],
        [
            "This is the jump from app bundle to reusable API surface.",
            "Developers can treat multiple modalities as one edge platform, not multiple unrelated demos.",
            "Routing and arbitration are platform behaviors customers can build on top of.",
            "If validated and hardened, this becomes Qualcomm's strongest developer-experience answer to Jetson.",
        ],
        C.GOLD,
        C.TEAL,
    )

    add_panel(slide, 7.1, 2.2, 5.52, 4.25, title="Suggested external wording", fill=C.PANEL_2)
    add_bullets(
        slide,
        7.3,
        2.72,
        5.1,
        3.3,
        [
            "Publicly visible today: containerized multimodal app package.",
            "Current internal milestone: unified OpenAI-style API surface across multiple modalities.",
            "Validated prototype: image understanding / vision-assistant flows through orchestrator.",
            "Next release objective: publish one canonical multimodal developer contract with onboarding and reference apps.",
        ],
        font_size=15,
        bullet_color=C.GOLD,
    )

    add_footer(
        slide,
        "Engineering refs (local): README.md; core-services/orchestrator/README.md; docs/API_CONTRACTS.md; core-services/image-to-text/README.md; core-services/speech-to-text/README.md; core-services/orchestrator/app/openai.py.",
        5,
    )
    add_notes(
        notes,
        5,
        "Internal capability",
        [
            "Be explicit that this is engineering evidence, not necessarily public GA.",
            "This is where the core platform opportunity becomes visible.",
            "The wording on the right side can be reused in leadership reviews to keep claims disciplined.",
        ],
    )


def slide_6_unified_api_value(prs: Presentation, notes: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        slide,
        "Why unified OpenAI-style multimodal endpoints matter",
        "Combining chat, transcription, speech generation, image generation, and vision flows is much more valuable than shipping isolated use-case demos.",
        ROADMAP_TAG,
    )

    layer_specs = [
        (0.95, 1.55, 11.45, 0.62, "App layer: copilots, operator UX, inspection assistants, voice interfaces, robotics workflows", C.PANEL_2),
        (1.25, 2.37, 10.85, 0.68, "Unified API layer: /v1/chat/completions | /v1/audio/* | /v1/images/* | multimodal routing", C.PANEL),
        (1.55, 3.28, 10.25, 0.76, "Modality services: text | STT | TTS | image generation | image understanding", C.PANEL_2),
        (1.85, 4.33, 9.65, 0.74, "Runtime and orchestration: routing, health, retries, arbitration, deployment packaging", C.PANEL),
        (2.15, 5.36, 9.05, 0.68, "IQ-9075 hardware foundation: local inference, industrial deployment, privacy-sensitive edge execution", C.PANEL_2),
    ]
    centers = []
    for x, y, w, h, text, fill in layer_specs:
        add_panel(slide, x, y, w, h, fill=fill)
        add_textbox(slide, x + 0.18, y + 0.19, w - 0.36, 0.3, text, font_size=15, color=C.WHITE, bold=True)
        centers.append((x + w / 2.0, y))
    for idx in range(len(centers) - 1):
        add_arrow(slide, centers[idx][0], centers[idx][1] + 0.62, centers[idx + 1][0], centers[idx + 1][1] - 0.02, C.CYAN)

    add_panel(slide, 9.56, 1.25, 3.05, 1.68, title="What changes for developers", fill=C.PANEL)
    add_bullets(
        slide,
        9.74,
        1.74,
        2.72,
        1.0,
        [
            "One client pattern.",
            "Less glue code.",
            "Composable app workflows.",
            "Faster migration from cloud-style prototypes to edge deployment.",
        ],
        font_size=13,
        bullet_color=C.TEAL,
    )

    add_footer(
        slide,
        "Public rationale supported by edge AI adoption sources. Internal route examples from local orchestrator docs/code should be validated before external use.",
        6,
    )
    add_notes(
        notes,
        6,
        "Unified API value",
        [
            "This is the strategic center of the deck.",
            "Endpoint unification matters because it converts many point capabilities into one developer platform.",
            "Keep the message focused on reduced friction and application composition.",
        ],
    )


def slide_7_developer_value(prs: Presentation, notes: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        slide,
        "Developer value: from demo bundle to reusable edge AI platform",
        "A platform unlock happens when the developer builds once against a consistent surface and then composes multiple workloads without re-learning each service.",
        ROADMAP_TAG,
    )

    rows = [
        ["Question", "Isolated demos", "Reusable edge platform"],
        ["Client integration", "Different commands and app-specific flows", "One API style and shared request patterns"],
        ["Composition", "Manual stitching across point solutions", "Chat + voice + image workflows can be composed"],
        ["Operations", "Per-service troubleshooting and ad hoc UX", "Gateway, health, status, routing, retries, policy"],
        ["Reuse", "One demo per use case", "Many customer apps on the same primitives"],
        ["Migration path", "Hard to port cloud-style prototypes", "OpenAI-style patterns lower developer switching cost"],
    ]
    add_matrix_table(slide, 0.62, 1.75, 12.0, 3.0, rows, [2.4, 4.1, 5.5])

    add_two_col_compare(
        slide,
        0.62,
        5.0,
        12.0,
        1.35,
        "What Jetson already gets right publicly",
        "What Qualcomm should emphasize",
        [
            "Developers see many official on-ramps: JetPack, Jetson AI tutorials, forums, Isaac ROS, cloud-to-edge stack.",
        ],
        [
            "Qualcomm should make GenAI Studio the familiar entry point that turns AI Engine Direct strength into developer velocity on IQ9-class hardware.",
        ],
        C.LIME,
        C.CYAN,
    )

    add_footer(
        slide,
        "Sources: NVIDIA Jetson software, getting-started, tutorials, and Isaac ROS pages; local GenAI Studio orchestrator documentation.",
        7,
    )
    add_notes(
        notes,
        7,
        "Developer value",
        [
            "Translate product strategy into developer ergonomics.",
            "This is where you make clear that the goal is not a prettier demo; it is lower integration friction and higher reuse.",
            "Use the last callout to bridge into the Jetson comparison.",
        ],
    )


def slide_8_customer_value(prs: Presentation, notes: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        slide,
        "Why customers care: one local stack can unlock many solution families",
        "Frame demand as customer pull signals: privacy, offline operation, lower latency, lower recurring cost, and easier integration into existing app stacks.",
        ROADMAP_TAG,
    )

    rows = [
        ["Use-case family", "Primary modalities", "Why local edge matters"],
        ["Operator copilot / HMI assistant", "Chat + TTS + STT", "Hands-free workflows, low latency, local site data"],
        ["Inspection and troubleshooting assistant", "Vision + chat + TTS", "Visual reasoning near equipment, limited connectivity"],
        ["Secure transcription / summarization", "STT + chat", "Data control and reduced cloud cost"],
        ["Field service or warehouse assistant", "Vision + voice + text", "Offline resilience and fast human interaction"],
        ["Creative / branded asset generation", "Image generation + chat", "Local iteration and IP-sensitive workflows"],
    ]
    add_matrix_table(slide, 0.62, 1.72, 7.1, 3.7, rows, [2.4, 2.2, 2.5])

    add_panel(slide, 7.98, 1.72, 4.64, 3.7, title="Value drivers customers recognize", fill=C.PANEL)
    add_bullets(
        slide,
        8.18,
        2.22,
        4.25,
        2.7,
        [
            "Low latency and usable UX without depending on cloud round trips.",
            "Better privacy, data residency, and local control over sensitive workflows.",
            "Reduced recurring inference cost versus always-on hosted APIs.",
            "One reusable integration model for many applications instead of stitching point solutions together.",
        ],
        font_size=15,
        bullet_color=C.TEAL,
    )

    add_panel(slide, 7.98, 5.58, 4.64, 0.74, fill=C.BG_2)
    add_textbox(
        slide,
        8.18,
        5.82,
        4.2,
        0.2,
        "Message discipline: say 'customer demand signals' unless validated customer evidence is available.",
        font_size=11,
        color=C.GOLD,
        bold=True,
    )

    add_footer(
        slide,
        "Customer-value framing supported by analyst summaries on local processing, privacy, and latency. Use-case examples are strategy illustrations, not customer commitments.",
        8,
    )
    add_notes(
        notes,
        8,
        "Customer value",
        [
            "Connect modality unification to customer outcomes, not just developer convenience.",
            "Keep examples realistic and industrially relevant.",
            "Avoid overclaiming customer validation if the evidence has not been formally collected.",
        ],
    )


def slide_9_jetson_mindshare(prs: Presentation, notes: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        slide,
        "Why Jetson wins mindshare today",
        "NVIDIA's edge advantage is not only silicon. It is public ecosystem packaging: price, tutorials, software stack, community, and robotics workflows.",
        PUBLIC_TAG,
    )

    add_stat_card(slide, 0.62, 1.48, 2.25, 0.95, "$249", "Jetson Orin Nano Super developer kit price", C.LIME)
    add_stat_card(slide, 3.02, 1.48, 2.25, 0.95, "67 TOPS", "Public Jetson Orin Nano Super claim", C.LIME)
    add_stat_card(slide, 5.42, 1.48, 2.55, 0.95, "JetPack + CUDA + TensorRT", "Integrated public software stack", C.LIME)
    add_stat_card(slide, 8.12, 1.48, 2.7, 0.95, "Isaac ROS + tutorials + forums", "Visible robotics and developer ecosystem", C.LIME)

    add_two_col_compare(
        slide,
        0.62,
        2.78,
        12.0,
        3.38,
        "Jetson public strengths",
        "Why these strengths matter",
        [
            "Affordable entry point marketed to developers, students, and makers.",
            "Extensive public stack: JetPack, CUDA, TensorRT, cuDNN, VPI, DLA.",
            "Official robotics surface via Isaac ROS and related workflows.",
            "Large learning funnel: getting-started guides, tutorials, forums, AI Lab content, blog posts.",
            "Frequent public software optimization story reinforces platform momentum.",
        ],
        [
            "Developers can start cheaply and see a clear path from tutorial to product prototype.",
            "The software stack is part of the brand, not hidden behind the hardware.",
            "Robotics teams get a ready-made ecosystem, not just a board.",
            "Public community visibility creates mindshare compounding effects.",
            "NVIDIA repeatedly turns software updates into visible platform wins.",
        ],
        C.LIME,
        C.WHITE,
    )

    add_footer(
        slide,
        "Sources: NVIDIA Jetson Orin Nano Super page; Jetson software page; JetPack 6.1 page; Isaac ROS page; Jetson getting-started and tutorials pages; NVIDIA Jetson edge AI blog.",
        9,
    )
    add_notes(
        notes,
        9,
        "Jetson mindshare",
        [
            "Be direct and fair: NVIDIA's current advantage is software ecosystem and public packaging.",
            "This slide should sound respectful, not defensive.",
            "Use it to build urgency for Qualcomm's platform-story response.",
        ],
    )


def slide_10_qualcomm_public_gaps(prs: Presentation, notes: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        slide,
        "Where Qualcomm is lagging publicly",
        "This is a public-positioning assessment, not a claim that Qualcomm lacks the underlying silicon or engineering potential.",
        INFERENCE_TAG,
    )

    add_panel(slide, 0.62, 1.55, 5.75, 4.8, title="Likely public gaps", fill=C.PANEL)
    add_bullets(
        slide,
        0.84,
        2.04,
        5.3,
        3.95,
        [
            "Lower public mindshare versus Jetson among edge AI developers.",
            "Public onboarding appears less streamlined and less visibly tutorial-rich.",
            "Public examples look more like separate use-case apps than a single reusable platform contract.",
            "Weaker public software-platform branding relative to JetPack / CUDA / Isaac-style narratives.",
            "Lower visible community momentum around end-to-end GenAI application building.",
        ],
        font_size=15,
        bullet_color=C.CORAL,
    )

    add_panel(slide, 6.66, 1.55, 5.96, 4.8, title="How to read this slide", fill=C.PANEL_2)
    add_bullets(
        slide,
        6.88,
        2.04,
        5.5,
        3.95,
        [
            "This is an inference from public source breadth reviewed on April 28, 2026.",
            "NVIDIA offers many visible on-ramps: price-point messaging, getting-started resources, forums, robotics stack, and public tutorials.",
            "Qualcomm public material is credible but narrower: strong hardware positioning plus a documented app bundle, with less obvious 'developer platform' packaging.",
            "That is a fixable positioning and productization gap, not a reason to retreat.",
        ],
        font_size=15,
        bullet_color=C.GOLD,
    )

    add_footer(
        slide,
        "Assessment based on public materials reviewed: Qualcomm GenAI Studio docs/repo, IQ-9075 / IQ9 pages, NVIDIA Jetson software, tutorials, forums, and robotics ecosystem pages.",
        10,
    )
    add_notes(
        notes,
        10,
        "Qualcomm public gaps",
        [
            "Keep the distinction clear: the lag is mostly public platform packaging and developer mindshare, not necessarily silicon capability.",
            "This slide should motivate action, not trigger defensiveness.",
            "Use the 'fixable gap' phrase explicitly when presenting.",
        ],
    )


def slide_11_gap_closure(prs: Presentation, notes: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        slide,
        "How this solution closes the gap",
        "The opportunity is to convert Qualcomm's hardware and runtime strength into a platform developers can recognize, adopt, and reuse.",
        ROADMAP_TAG,
    )

    rows = [
        ["Current gap", "Response with GenAI Studio platform direction", "Strategic effect"],
        ["Fragmented public app perception", "Publish one canonical multimodal platform narrative", "Shifts story from demos to reusable developer surface"],
        ["Weak public API familiarity", "Expose OpenAI-style routes consistently through orchestrator", "Lowers developer switching cost from cloud patterns"],
        ["Higher integration friction", "Unify routing, health, retries, and modality composition", "Accelerates prototyping and productization"],
        ["Jetson ecosystem narrative advantage", "Ship reference apps, tutorials, and migration paths on IQ-9075", "Turns silicon value into ecosystem momentum"],
        ["Unclear industrial AI differentiation", "Tie the platform directly to rugged IQ9 deployment benefits", "Positions Qualcomm away from hobbyist-only comparison traps"],
    ]
    add_matrix_table(slide, 0.62, 1.72, 12.0, 4.05, rows, [2.2, 5.1, 4.7])

    add_panel(slide, 0.62, 5.97, 12.0, 0.6, fill=C.BG_2)
    add_textbox(
        slide,
        0.82,
        6.15,
        11.55,
        0.22,
        "Executive message: this is not about copying Jetson. It is about making Qualcomm easier to build on, more modular to deploy, and more differentiated in industrial edge AI.",
        font_size=13,
        color=C.WHITE,
        bold=True,
    )

    add_footer(
        slide,
        "Strategy synthesis based on public sources plus current local engineering docs. External claims should remain within validated public or owner-approved bounds.",
        11,
    )
    add_notes(
        notes,
        11,
        "Closing the gap",
        [
            "This is the answer slide. Every previous slide feeds into this table.",
            "Keep the message focused on platform conversion, familiar APIs, and industrial differentiation.",
            "The closing bar sentence is the headline takeaway for leadership.",
        ],
    )


def slide_12_compare_matrix(prs: Presentation, notes: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        slide,
        "Competitive differentiation: Qualcomm vs Jetson",
        "Jetson is ahead in public ecosystem maturity. Qualcomm can differentiate around industrial-grade local multimodal platformization if it makes the software story visible.",
        ROADMAP_TAG,
    )

    rows = [
        ["Dimension", "Qualcomm GenAI Studio on IQ-9075", "Jetson Orin Nano / Jetson stack"],
        ["Public positioning", "Industrial-grade AI platform emerging", "Developer-first edge AI ecosystem established"],
        ["Entry experience", "Stronger if simplified through one documented platform path", "Very strong public on-ramp via price, guides, tutorials, community"],
        ["Industrial deployment fit", "Strong public hardware story: rugged, longevity, Ubuntu + Qualcomm Linux", "Strong prototyping story; industrial fit varies by module / system partner"],
        ["OpenAI-style API familiarity", "High potential; current internal direction appears strong", "Can run popular open-model stacks, but public identity centers more on JetPack/CUDA ecosystem than one OpenAI-compatible product layer"],
        ["Robotics ecosystem", "Developing public story", "Very strong via Isaac ROS and related workflows"],
        ["Multimodal local platform story", "Potentially differentiating if made public and productized", "Strong model breadth and tutorials; platform story already visible"],
        ["Best strategic lane", "Industrial edge, privacy-sensitive local AI, reusable multimodal endpoints", "Broad developer mindshare, robotics, education, and ecosystem breadth"],
    ]
    add_matrix_table(slide, 0.54, 1.6, 12.24, 4.95, rows, [1.8, 5.1, 5.34])

    add_footer(
        slide,
        "Sources: Qualcomm IQ-9075 / IQ9 / GenAI Studio public materials; NVIDIA Jetson Orin Nano Super, JetPack, Jetson software, Isaac ROS, and tutorials pages.",
        12,
    )
    add_notes(
        notes,
        12,
        "Competitive differentiation matrix",
        [
            "Keep this balanced. The credibility comes from honesty about Jetson's strength.",
            "Highlight Qualcomm's best lane: industrial-grade, local, multimodal, reusable platform story.",
            "Do not overclaim that Qualcomm already wins on software ecosystem breadth.",
        ],
    )


def slide_13_gtm(prs: Presentation, notes: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        slide,
        "Go-to-market and ecosystem implications",
        "If GenAI Studio is treated as a platform, not a demo bundle, Qualcomm can create a clearer adoption funnel for OEMs, partners, and developers.",
        ROADMAP_TAG,
    )

    add_panel(slide, 0.62, 1.62, 3.8, 4.92, title="Messaging shift", fill=C.PANEL)
    add_bullets(
        slide,
        0.84,
        2.12,
        3.35,
        4.1,
        [
            "From: 'four AI apps on IQ9'",
            "To: 'OpenAI-compatible edge AI platform for industrial multimodal apps'",
            "Lead with reusable endpoints, on-device privacy, low latency, and rugged deployment fit.",
        ],
        font_size=16,
        bullet_color=C.CYAN,
    )

    add_panel(slide, 4.65, 1.62, 3.95, 4.92, title="Ecosystem motion", fill=C.PANEL)
    add_bullets(
        slide,
        4.87,
        2.12,
        3.45,
        4.1,
        [
            "Publish canonical quickstart and one-click developer journey.",
            "Ship 2-3 reference solutions: operator copilot, secure transcription, visual inspection assistant.",
            "Create partner enablement kits for robotics and industrial OEM channels.",
            "Tie GenAI Studio to broader Qualcomm developer tooling and AI Engine Direct narrative.",
        ],
        font_size=15,
        bullet_color=C.TEAL,
    )

    add_panel(slide, 8.83, 1.62, 3.79, 4.92, title="Commercial effect", fill=C.PANEL)
    add_bullets(
        slide,
        9.05,
        2.12,
        3.35,
        4.1,
        [
            "Better story for design-ins where customers want reusable local AI building blocks.",
            "Stronger defense against 'Jetson is easier to start with'.",
            "Higher perceived software maturity without needing to match NVIDIA tutorial-for-tutorial on day one.",
            "Better leverage of IQ9 industrial differentiation in customer conversations.",
        ],
        font_size=15,
        bullet_color=C.GOLD,
    )

    add_footer(
        slide,
        "Strategy implications inferred from reviewed public positioning and current local platform direction.",
        13,
    )
    add_notes(
        notes,
        13,
        "GTM and ecosystem implications",
        [
            "Translate platform design into go-to-market motion.",
            "The core change is messaging plus reference solutions plus simpler onboarding.",
            "Use this slide to connect platform work to commercial leverage.",
        ],
    )


def slide_14_priorities(prs: Presentation, notes: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        slide,
        "Recommended next actions and product priorities",
        "The fastest way to close the gap is to productize the software story around one canonical developer surface.",
        ROADMAP_TAG,
    )

    steps = [
        ("1", "Public product narrative", "Reposition GenAI Studio publicly as an edge AI platform starter kit, not only four separate applications."),
        ("2", "Canonical API contract", "Publish and stabilize one external OpenAI-style developer contract with clear GA vs preview labels."),
        ("3", "One-click onboarding", "Reduce setup friction on Ubuntu IQ-9075 with a single happy-path quickstart and validation tool."),
        ("4", "Reference apps", "Publish 2-3 polished reference solutions that combine modalities around real industrial workflows."),
        ("5", "Proof and benchmarks", "Add platform-level proof: latency, privacy, offline, and deployment value, not only raw model execution."),
        ("6", "Jetson response kit", "Create competitive messaging, migration notes, and partner enablement material for field teams."),
    ]

    start_y = 1.62
    for idx, (num, title, body) in enumerate(steps):
        y = start_y + idx * 0.77
        circ = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.78), Inches(y), Inches(0.46), Inches(0.46))
        circ.fill.solid()
        circ.fill.fore_color.rgb = [C.CYAN, C.TEAL, C.GOLD, C.CORAL, C.PURPLE, C.LIME][idx]
        circ.line.fill.background()
        ctf = circ.text_frame
        ctf.clear()
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = num
        cr.font.size = Pt(14)
        cr.font.bold = True
        cr.font.color.rgb = C.BG
        cr.font.name = "Aptos"

        add_panel(slide, 1.42, y - 0.05, 11.1, 0.58, fill=C.PANEL if idx % 2 == 0 else C.PANEL_2)
        add_textbox(slide, 1.62, y + 0.02, 2.35, 0.18, title, font_size=14, color=C.WHITE, bold=True)
        add_textbox(slide, 3.95, y + 0.01, 8.25, 0.22, body, font_size=12, color=C.MUTED)

    add_panel(slide, 0.78, 6.4, 11.75, 0.38, fill=C.BG_2)
    add_textbox(
        slide,
        0.96,
        6.49,
        11.3,
        0.15,
        "Executive ask: align product, engineering, and GTM on one validated public platform narrative and one canonical developer journey.",
        font_size=12,
        color=C.WHITE,
        bold=True,
    )

    add_footer(
        slide,
        "Roadmap recommendations synthesized from source review and current local engineering direction.",
        14,
    )
    add_notes(
        notes,
        14,
        "Priorities and ask",
        [
            "Make the action list concrete and platform-oriented.",
            "The ask is cross-functional alignment: product, engineering, docs, and GTM must tell the same story.",
            "This slide should be easy for leadership to turn into owners and milestones.",
        ],
    )


def slide_15_appendix(prs: Presentation, notes: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        slide,
        "Appendix: source map and validation boundaries",
        "Use this slide to preserve accuracy discipline when the deck is reused or externalized.",
        PUBLIC_TAG,
    )

    add_panel(slide, 0.62, 1.52, 4.0, 4.85, title="Public Qualcomm sources", fill=C.PANEL)
    add_textbox(
        slide,
        0.82,
        2.0,
        3.6,
        4.1,
        "1. GenAI Studio doc (updated Feb 23, 2026)\n"
        "2. Public GenAI Studio GitHub README\n"
        "3. IQ-9075 product page\n"
        "4. IQ-9075 EVK developer page\n"
        "5. IQ-9075 EVK quickstart / documentation\n"
        "6. IQ9 public product brief / docs\n",
        font_size=13,
        color=C.WHITE,
    )

    add_panel(slide, 4.85, 1.52, 4.0, 4.85, title="Public NVIDIA and market sources", fill=C.PANEL)
    add_textbox(
        slide,
        5.05,
        2.0,
        3.6,
        4.1,
        "1. Jetson Orin Nano Super page\n"
        "2. Jetson software page\n"
        "3. JetPack 6.1 page\n"
        "4. Isaac ROS page\n"
        "5. Jetson getting-started / tutorials\n"
        "6. Jetson edge AI blog\n"
        "7. Global Market Insights edge AI market\n"
        "8. MarketsandMarkets edge AI hardware\n"
        "9. EDGE AI Foundation 2025 report\n",
        font_size=13,
        color=C.WHITE,
    )

    add_panel(slide, 9.08, 1.52, 3.54, 4.85, title="Internal validation boundary", fill=C.PANEL_2)
    add_bullets(
        slide,
        9.28,
        2.0,
        3.1,
        3.95,
        [
            "Local repo evidence should be labeled 'engineering evidence' until owner-validated.",
            "Do not present OpenAI-style multimodal unification as public GA unless official Qualcomm sources confirm it.",
            "Do not use unsupported market-share claims.",
            "Preserve exact dates when discussing what is public 'today'.",
        ],
        font_size=13,
        bullet_color=C.GOLD,
    )

    add_footer(
        slide,
        "Deck date basis: April 28, 2026. See accompanying notes file for slide-by-slide presenter guidance.",
        15,
    )
    add_notes(
        notes,
        15,
        "Appendix",
        [
            "This slide protects the deck from being misused later.",
            "Public vs internal vs roadmap separation is the most important accuracy control in the whole deck.",
            "Use the notes file if the presentation needs a spoken narrative.",
        ],
    )


if __name__ == "__main__":
    build_deck()
    print(f"Wrote {OUT_PPTX}")
    print(f"Wrote {OUT_NOTES}")
